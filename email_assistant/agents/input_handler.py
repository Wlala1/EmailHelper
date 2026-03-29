from __future__ import annotations

import base64
import json
import mimetypes
import os
import re
import tempfile
import zipfile
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from config import MAX_ATTACHMENT_TEXT_LENGTH, OPENAI_API_KEY, OPENAI_BASE_URL, OPENAI_MODEL

try:
    import markdownify
except Exception:  # pragma: no cover - optional dependency
    markdownify = None

try:
    import mammoth
except Exception:  # pragma: no cover - optional dependency
    mammoth = None

try:
    import openpyxl
except Exception:  # pragma: no cover - optional dependency
    openpyxl = None

try:
    import pdfminer.high_level
except Exception:  # pragma: no cover - optional dependency
    pdfminer = None
else:
    import pdfminer

try:
    import pptx
except Exception:  # pragma: no cover - optional dependency
    pptx = None

try:
    from bs4 import BeautifulSoup
except Exception:  # pragma: no cover - optional dependency
    BeautifulSoup = None

try:
    from markitdown import MarkItDown
except Exception:  # pragma: no cover - optional dependency
    MarkItDown = None

try:
    from openai import OpenAI
except Exception:  # pragma: no cover - optional dependency
    OpenAI = None

IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "gif", "webp"}
AUDIO_EXTENSIONS = {"wav", "mp3", "m4a"}
VIDEO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "webm"}
TEXT_EXTENSIONS = {"txt", "md", "sh", "yaml", "yml", "toml", "csv", "py"}
JSON_EXTENSIONS = {"json", "jsonld"}

DATE_PATTERNS = [
    r"\b\d{4}-\d{2}-\d{2}\b",
    r"\b\d{4}/\d{2}/\d{2}\b",
    r"\b\d{1,2}/\d{1,2}/\d{4}\b",
    r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{4}\b",
]


@dataclass
class ParsedAttachment:
    attachment_id: str
    doc_type: Optional[str]
    relevance_score: float
    topics: list[str]
    named_entities: list[str]
    time_expressions: list[str]
    extracted_text: str


def _truncate(text: str, limit: int = MAX_ATTACHMENT_TEXT_LENGTH) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + "\n... [truncated]"


def _extract_keywords(text: str, limit: int = 8) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z0-9_]{3,}", text.lower())
    stopwords = {
        "this",
        "that",
        "with",
        "from",
        "your",
        "have",
        "will",
        "please",
        "subject",
        "email",
        "about",
        "thanks",
        "dear",
        "team",
        "regards",
        "conference",
    }
    words = [w for w in words if w not in stopwords]
    counts = Counter(words)
    return [k for k, _ in counts.most_common(limit)]


def _extract_named_entities(text: str, sender_email: Optional[str] = None) -> list[str]:
    entities: set[str] = set()
    for email in re.findall(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", text):
        entities.add(email)
    if sender_email:
        entities.add(sender_email)

    for match in re.findall(r"\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b", text):
        entities.add(match)
    return sorted(entities)[:30]


def _extract_time_expressions(text: str) -> list[str]:
    found: list[str] = []
    for pattern in DATE_PATTERNS:
        found.extend(re.findall(pattern, text))
    # Keep insertion order while removing duplicates.
    return list(dict.fromkeys(found))[:20]


def _infer_doc_type(filename: str, text: str) -> Optional[str]:
    lower = f"{filename}\n{text}".lower()
    if "cfp" in lower or "call for papers" in lower:
        return "cfp"
    if "meeting" in lower or "agenda" in lower:
        return "meeting"
    if "invoice" in lower or "receipt" in lower:
        return "invoice"
    if "course" in lower or "canvas" in lower:
        return "course_update"
    if "career" in lower or "job" in lower or "intern" in lower:
        return "career"
    if "event" in lower or "social" in lower:
        return "social_event"
    return None


def _guess_content_type(path: str, fallback: Optional[str]) -> str:
    if fallback:
        return fallback
    mime, _ = mimetypes.guess_type(path)
    return mime or "application/octet-stream"


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_json(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return json.dumps(json.load(f), ensure_ascii=False, indent=2)


def _read_pdf(path: str) -> str:
    if pdfminer is None:
        raise RuntimeError("pdfminer is not installed")
    return pdfminer.high_level.extract_text(path)


def _read_docx(path: str) -> str:
    if mammoth is not None:
        with open(path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
        if BeautifulSoup is not None:
            soup = BeautifulSoup(result.value, "html.parser")
            return soup.get_text("\n")
        return result.value

    if MarkItDown is not None:
        md = MarkItDown(enable_plugins=True)
        return md.convert(path).text_content

    raise RuntimeError("mammoth/markitdown not available for docx parsing")


def _read_html(path: str) -> str:
    raw = _read_text(path)
    if BeautifulSoup is not None:
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style"]):
            tag.extract()
        if markdownify is not None:
            return markdownify.markdownify(str(soup), heading_style="ATX")
        return soup.get_text("\n")
    return raw


def _read_pptx(path: str) -> str:
    if pptx is None:
        raise RuntimeError("python-pptx is not installed")

    prs = pptx.Presentation(path)
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        chunks.append(f"Slide {idx}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    chunks.append(text)
    return "\n".join(chunks)


def _read_xlsx(path: str) -> str:
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed")

    wb = openpyxl.load_workbook(path, data_only=True)
    chunks: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        chunks.append(f"Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v) if v is not None else "" for v in row]
            if any(values):
                chunks.append("\t".join(values))
    return "\n".join(chunks)


def _caption_with_openai(path: str, media_type: str) -> str:
    if OpenAI is None or not OPENAI_API_KEY:
        return f"[{media_type} caption unavailable: OPENAI_API_KEY not set]"

    client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL)
    if media_type == "audio":
        with open(path, "rb") as audio_file:
            result = client.audio.transcriptions.create(model="gpt-4o-transcribe", file=audio_file)
        return result.text or "[audio transcription empty]"

    prompt = {
        "image": "Describe the image in detail and capture text/deadlines/topics if visible.",
        "video": "Describe the video and list any visible text, dates, or action items.",
    }[media_type]
    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")

    ext = Path(path).suffix.lower()
    mime = mimetypes.types_map.get(ext, "application/octet-stream")
    response = client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }
        ],
        temperature=0,
    )
    return response.choices[0].message.content or f"[{media_type} caption empty]"


def _read_zip(path: str) -> str:
    chunks: list[str] = [f"ZIP archive: {Path(path).name}"]
    with tempfile.TemporaryDirectory(prefix="ouma_zip_") as temp_dir:
        with zipfile.ZipFile(path, "r") as zf:
            zf.extractall(temp_dir)

        for root, _, files in os.walk(temp_dir):
            for name in files:
                local = Path(root) / name
                rel = local.relative_to(temp_dir)
                try:
                    nested = parse_attachment_content(str(local), name=name, content_type=None)
                    chunks.append(f"File: {rel}\n{nested}")
                except Exception as exc:
                    chunks.append(f"File: {rel}\n[failed to parse: {exc}]")
    return "\n\n".join(chunks)


def _fallback_with_markitdown(path: str) -> Optional[str]:
    if MarkItDown is None:
        return None
    try:
        md = MarkItDown(enable_plugins=True)
        result = md.convert(path)
        return getattr(result, "text_content", None)
    except Exception:
        return None


def parse_attachment_content(path: str, name: Optional[str], content_type: Optional[str]) -> str:
    ext = Path(path).suffix.lower().lstrip(".")

    if ext in TEXT_EXTENSIONS:
        return _read_text(path)
    if ext in JSON_EXTENSIONS:
        return _read_json(path)
    if ext == "pdf":
        return _read_pdf(path)
    if ext in {"docx", "doc"}:
        return _read_docx(path)
    if ext in {"html", "htm"}:
        return _read_html(path)
    if ext in {"pptx", "ppt"}:
        return _read_pptx(path)
    if ext in {"xlsx", "xls"}:
        return _read_xlsx(path)
    if ext == "zip":
        return _read_zip(path)
    if ext in IMAGE_EXTENSIONS:
        return _caption_with_openai(path, "image")
    if ext in AUDIO_EXTENSIONS:
        return _caption_with_openai(path, "audio")
    if ext in VIDEO_EXTENSIONS:
        return _caption_with_openai(path, "video")

    fallback = _fallback_with_markitdown(path)
    if fallback:
        return fallback
    return f"[unsupported attachment type: {ext or _guess_content_type(path, content_type)}]"


def parse_attachment(
    *,
    attachment_id: str,
    name: str,
    path: str,
    content_type: Optional[str],
    sender_email: Optional[str] = None,
) -> ParsedAttachment:
    raw = parse_attachment_content(path, name=name, content_type=content_type)
    extracted_text = _truncate(raw.strip())
    topics = _extract_keywords(extracted_text)
    named_entities = _extract_named_entities(extracted_text, sender_email=sender_email)
    time_expressions = _extract_time_expressions(extracted_text)
    doc_type = _infer_doc_type(name, extracted_text)

    # Simple relevance score: mixed signal from text richness and whether dates/entities were extracted.
    score = 0.2
    if len(extracted_text) > 200:
        score += 0.3
    if topics:
        score += 0.2
    if named_entities:
        score += 0.2
    if time_expressions:
        score += 0.1

    return ParsedAttachment(
        attachment_id=attachment_id,
        doc_type=doc_type,
        relevance_score=min(1.0, round(score, 4)),
        topics=topics,
        named_entities=named_entities,
        time_expressions=time_expressions,
        extracted_text=extracted_text,
    )
